"""PowerPoint reading and export helpers."""

from pathlib import Path

from pptx import Presentation

from . import runtime as cfg

def get_slide_notes(pptx_path: str) -> list[str]:
    """Read speaker notes from each slide in the PPTX file."""
    prs = Presentation(pptx_path)
    notes = []
    for i, slide in enumerate(prs.slides):
        text = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            text = tf.text.strip()
        notes.append(text)
        status = f"({len(text)} chars)" if text else "(empty)"
        print(f"   Slide {i+1}: {status}")
    return notes


def export_slides_as_images(pptx_path: str, output_dir: Path) -> list[str]:
    """Export each slide as a PNG through PowerPoint COM."""
    import win32com.client

    abs_path = str(Path(pptx_path).resolve())
    pp = win32com.client.Dispatch("Powerpoint.Application")
    pp.Visible = True  # PowerPoint must be visible for slide export.

    try:
        prs = pp.Presentations.Open(abs_path, ReadOnly=True)
        paths = []
        for i in range(1, prs.Slides.Count + 1):
            img_path = str(output_dir.resolve() / f"slide_{i:03d}.png")
            prs.Slides(i).Export(img_path, "PNG", cfg.SLIDE_WIDTH, cfg.SLIDE_HEIGHT)
            paths.append(img_path)
            print(f"   Slide {i}/{prs.Slides.Count}")
        prs.Close()
        return paths
    finally:
        pp.Quit()


def get_gif_overlays(pptx_path: str, slide_index: int, output_dir: Path) -> list[dict]:
    """Extract GIF images from a slide and return their paths and positions."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]

    scale_x = cfg.SLIDE_WIDTH / prs.slide_width
    scale_y = cfg.SLIDE_HEIGHT / prs.slide_height

    gifs = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            img = shape.image
            if img.content_type == "image/gif":
                gif_path = str(output_dir / f"slide_{slide_index+1:03d}_gif{len(gifs)}.gif")
                with open(gif_path, "wb") as f:
                    f.write(img.blob)
                gifs.append({
                    "path": gif_path,
                    "left":   int(shape.left   * scale_x),
                    "top":    int(shape.top    * scale_y),
                    "width":  int(shape.width  * scale_x),
                    "height": int(shape.height * scale_y),
                })
    return gifs


